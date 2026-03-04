#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Generate figures (and optional PPT) for the institutional 80/20 split.

Inputs:
  - split_summary_patients.xlsx  (produced by the split script)

Outputs:
  - PNG figures (age, gender, ethnicity, affected, raw Waldenström, IIb/IIIa subset lateral pillar, split sizes)
  - Optional PPT: split_methods_demographics.pptx

How to run (example):
  python make_split_figures.py
"""

import sys
from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.ticker import PercentFormatter

# Optional: PowerPoint; if missing, the script still generates PNGs.
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPT_AVAILABLE = True
except Exception:
    PPT_AVAILABLE = False


# ---------------------------
# 🧭 Paths to set (EDIT ME)
# ---------------------------
# 1) Location of the Excel report created by the split script:
REPORT_XLSX = Path(r"split_summary_patients.xlsx")
#    If it’s in OneDrive, set it explicitly, e.g.:
# REPORT_XLSX = Path(r"C:\Users\SR207348\OneDrive - Scottish Rite for Children\Kim Research\HIPMETRICS-lateral pillar\demographic split\split_summary_patients.xlsx")

# 2) Output folder for generated figures/PPT (your OneDrive folder is fine):
OUT_DIR = Path(r"C:\Users\SR207348\OneDrive - Scottish Rite for Children\Kim Research\HIPMETRICS-lateral pillar\demographic split\figures_outputs")
#    You can change to any folder you prefer.


# ---------------------------
# Helpers
# ---------------------------
def ensure_exists(p: Path):
    p.mkdir(parents=True, exist_ok=True)
    return p

def save_side_by_side_bar(data_a, data_b, cats, labels, title, out_path, rotate=45):
    """
    Plot normalized (proportion) side-by-side bars for two splits.
    data_a / data_b: pandas Series indexed by category with counts (not proportions).
    """
    idx = np.arange(len(cats))
    A = np.array([float(data_a.get(c, 0.0)) for c in cats])
    B = np.array([float(data_b.get(c, 0.0)) for c in cats])

    A_prop = A / A.sum() if A.sum() > 0 else A
    B_prop = B / B.sum() if B.sum() > 0 else B

    plt.figure(figsize=(max(6, 0.5 * len(cats) + 2), 4))
    w = 0.4
    plt.bar(idx - w / 2, A_prop, width=w, label=labels[0])
    plt.bar(idx + w / 2, B_prop, width=w, label=labels[1])
    plt.xticks(idx, cats, rotation=rotate, ha="right")
    plt.gca().yaxis.set_major_formatter(PercentFormatter(1.0))
    plt.ylabel("Proportion")
    plt.title(title)
    plt.legend()
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()
    return out_path

def add_single_picture_slide(prs: "Presentation", title: str, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.3), width=Inches(12.3))


# ---------------------------
# Main
# ---------------------------
def main():
    if not REPORT_XLSX.exists():
        print(f"[ERROR] Report not found: {REPORT_XLSX}")
        print("Edit REPORT_XLSX at the top of the script to point to your split_summary_patients.xlsx.")
        sys.exit(1)

    out_dir = ensure_exists(OUT_DIR)

    # ---- Read needed sheets from the report
    try:
        summary = pd.read_excel(REPORT_XLSX, sheet_name="00_summary")
        dists   = pd.read_excel(REPORT_XLSX, sheet_name="04_distributions_patients")
        assign  = pd.read_excel(REPORT_XLSX, sheet_name="05_split_assignment_institutions")
    except ValueError as e:
        print(f"[ERROR] Could not read required sheets from {REPORT_XLSX}: {e}")
        print("Make sure this is the Excel produced by the split script.")
        sys.exit(1)

    # Optional sheets (if present)
    try:
        wald_raw = pd.read_excel(REPORT_XLSX, sheet_name="08_waldenstrom_stage_raw")
    except Exception:
        wald_raw = None
    try:
        subset_lp = pd.read_excel(REPORT_XLSX, sheet_name="09_subset_Wald_IIb_IIIa_with_lateral_pillar")
    except Exception:
        subset_lp = None

    # ---- 1) Total patients by split (counts)
    split_counts = assign.groupby("split")["n_patients"].sum().to_dict()
    plt.figure(figsize=(6, 4))
    plt.bar(["Train/Val", "Test"],
            [split_counts.get("Train/Val", 0), split_counts.get("Test", 0)],
            color=["#4e79a7", "#f28e2c"])
    plt.title("Total patients by split")
    plt.ylabel("Patients")
    plt.tight_layout()
    fp_sizes = out_dir / "split_sizes.png"
    plt.savefig(fp_sizes, dpi=200)
    plt.close()

    # ---- Helper to extract two series of counts for a given variable
    def two_series(var_name: str):
        df = dists[dists["variable"] == var_name]
        cats = sorted(set(df["category"]))
        tr = df[df["split"] == "Train/Val"].set_index("category")["count"]
        te = df[df["split"] == "Test"].set_index("category")["count"]
        return tr, te, cats

    # ---- 2) Age
    tr_age, te_age, cats_age = two_series("age")
    fp_age = save_side_by_side_bar(tr_age, te_age, cats_age,
                                   ["Train/Val", "Test"],
                                   "Age distribution (patients)",
                                   out_dir / "age_distribution.png",
                                   rotate=45)

    # ---- 3) Gender
    tr_g, te_g, cats_g = two_series("gender")
    fp_gender = save_side_by_side_bar(tr_g, te_g, cats_g,
                                      ["Train/Val", "Test"],
                                      "Gender distribution (patients)",
                                      out_dir / "gender_distribution.png",
                                      rotate=0)

    # ---- 4) Race/Ethnicity
    tr_r, te_r, cats_r = two_series("race")
    fp_race = save_side_by_side_bar(tr_r, te_r, cats_r,
                                    ["Train/Val", "Test"],
                                    "Ethnicity distribution (patients)",
                                    out_dir / "ethnicity_distribution.png",
                                    rotate=45)

    # ---- 5) Affected vs Unaffected
    tr_a, te_a, cats_a = two_series("affected")
    fp_aff = save_side_by_side_bar(tr_a, te_a, cats_a,
                                   ["Train/Val", "Test"],
                                   "Affected vs Unaffected (patients)",
                                   out_dir / "affected_distribution.png",
                                   rotate=0)

    # ---- 6) Raw Waldenström stage (if available)
    fp_wraw = None
    if wald_raw is not None and not wald_raw.empty:
        wr_tr = wald_raw[wald_raw["split"] == "Train/Val"].set_index("waldenstrom_stage_raw")["count"]
        wr_te = wald_raw[wald_raw["split"] == "Test"].set_index("waldenstrom_stage_raw")["count"]
        cats_wr = sorted(set(wald_raw["waldenstrom_stage_raw"]))
        fp_wraw = save_side_by_side_bar(wr_tr, wr_te, cats_wr,
                                        ["Train/Val", "Test"],
                                        "Waldenström stage (raw)",
                                        out_dir / "waldenstrom_raw.png",
                                        rotate=0)

    # ---- 7) Subset: Waldenström IIb/IIIa + lateral pillar (if available)
    fp_subset = None
    if subset_lp is not None and not subset_lp.empty:
        tmp = subset_lp[subset_lp["split"].isin(["Train/Val", "Test"])]
        tmp = tmp[tmp["lateral_pillar"].astype(str).str.lower() != "chi2"]  # drop any STATS line
        tr_lp = tmp[tmp["split"] == "Train/Val"].set_index("lateral_pillar")["count"]
        te_lp = tmp[tmp["split"] == "Test"].set_index("lateral_pillar")["count"]
        cats_lp = sorted(set(tmp["lateral_pillar"]))
        fp_subset = save_side_by_side_bar(tr_lp, te_lp, cats_lp,
                                          ["Train/Val", "Test"],
                                          "Lateral pillar — Waldenström IIb/IIIa subset",
                                          out_dir / "subset_lateral_pillar.png",
                                          rotate=0)

    print(f"[OK] Saved figures to: {out_dir}")

    # ---- Optional: Build a PowerPoint
    if PPT_AVAILABLE:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Institutional 80/20 Split — Methods & Demographics"
        slide.placeholders[1].text = f"Auto-generated from: {REPORT_XLSX.name}"

        # Key metrics
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Split summary (key metrics)"
        sum_map = dict(zip(summary["metric"], summary["value"]))
        lines = [
            f"Total patients: {int(sum_map.get('total_unique_patients', 0))}",
            f"Train/Val patients: {int(sum_map.get('train_val_patients_n', 0))}",
            f"Test patients: {int(sum_map.get('test_patients_n', 0))}",
            f"Test fraction (patients): {float(sum_map.get('test_fraction_patients', 0.0)):.3f}",
            f"Test institutions: {int(sum_map.get('test_institutions_n', 0))} / "
            f"{int(sum_map.get('total_institutions_n', 0))} "
            f"({float(sum_map.get('test_fraction_institutions', 0.0)):.3f})",
        ]
        ph = slide.placeholders[1]
        ph.text = lines[0]
        for L in lines[1:]:
            p = ph.text_frame.add_paragraph()
            p.text = L

        # Add each figure as a slide
        add_single_picture_slide(prs, "Total patients by split", fp_sizes)
        add_single_picture_slide(prs, "Age distribution (patients)", fp_age)
        add_single_picture_slide(prs, "Gender distribution (patients)", fp_gender)
        add_single_picture_slide(prs, "Ethnicity distribution (patients)", fp_race)
        add_single_picture_slide(prs, "Affected vs Unaffected (patients)", fp_aff)
        if fp_wraw is not None:
            add_single_picture_slide(prs, "Waldenström stage (raw)", fp_wraw)
        if fp_subset is not None:
            add_single_picture_slide(prs, "Lateral pillar — Waldenström IIb/IIIa subset", fp_subset)

        ppt_path = out_dir / "split_methods_demographics.pptx"
        prs.save(ppt_path)
        print(f"[OK] PowerPoint saved to: {ppt_path}")
    else:
        print("[INFO] python-pptx not installed — skipped PPT. To include it, install: pip install python-pptx")


if __name__ == "__main__":
    main()